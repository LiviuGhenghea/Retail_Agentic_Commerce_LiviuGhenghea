"""Generate editable PowerPoint pitch deck for K9 Agent."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ───────────────────────────────────────────────────────────────────
ZURICH_BLUE  = RGBColor(0x00, 0x38, 0x82)
MID_BLUE     = RGBColor(0x00, 0x57, 0xB8)
LIGHT_BLUE   = RGBColor(0xE8, 0xF0, 0xFB)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MID_GRAY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF5, 0xF7, 0xFA)
GREEN        = RGBColor(0x1A, 0x7A, 0x4A)
LIGHT_GREEN  = RGBColor(0xE8, 0xF5, 0xEE)
ORANGE       = RGBColor(0xE0, 0x7B, 0x00)
ACCENT_BLUE  = RGBColor(0x4A, 0x90, 0xD9)

OUTPUT = "/Users/LIVIU.GHENGHEA/k9-agent/LiviuGhenghea_Retail_Agentic_Commerce.pptx"

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def inches(x): return Inches(x)
def pts(x): return Pt(x)

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=11, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = pts(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_label_box(slide, x, y, w, h, label, body, label_size=8, body_size=9,
                  bg=LIGHT_BLUE, label_color=MID_BLUE):
    """Coloured header card with label + body text."""
    # Header bar
    hdr_h = inches(0.28)
    hdr = add_rect(slide, x, y, w, hdr_h, fill_color=bg)
    add_textbox(slide, x + inches(0.08), y + inches(0.02), w - inches(0.16), hdr_h,
                label, font_size=label_size, bold=True, color=label_color)
    # Body
    body_h = h - hdr_h
    body_bg = add_rect(slide, x, y + hdr_h, w, body_h, fill_color=WHITE, line_color=MID_BLUE, line_width=Pt(0.5))
    add_textbox(slide, x + inches(0.1), y + hdr_h + inches(0.05), w - inches(0.2), body_h - inches(0.1),
                body, font_size=body_size, color=DARK_GRAY, wrap=True)

def add_section_label(slide, x, y, w, text):
    add_textbox(slide, x, y, w, inches(0.25), text,
                font_size=8, bold=True, color=MID_BLUE)

def add_zurich_logo(slide, x, y, w=inches(1.5), h=inches(0.55)):
    """Draw Zurich logo as blue rect + ZURICH text."""
    box = add_rect(slide, x, y, w, h, fill_color=ZURICH_BLUE)
    add_textbox(slide, x, y + inches(0.06), w, h - inches(0.06),
                "ZURICH", font_size=16, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, font_name="Calibri")

def slide_header(slide, num, title, subtitle):
    """Blue header bar with slide number, title, subtitle and Zurich logo."""
    bar_h = inches(0.85)
    add_rect(slide, 0, 0, W, bar_h, fill_color=ZURICH_BLUE)
    # Slide number
    add_textbox(slide, inches(0.2), inches(0.05), inches(0.7), bar_h,
                str(num), font_size=36, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    # Title + subtitle
    add_textbox(slide, inches(1.0), inches(0.08), inches(9.5), inches(0.45),
                title, font_size=20, bold=True, color=WHITE)
    add_textbox(slide, inches(1.0), inches(0.52), inches(9.5), inches(0.28),
                subtitle, font_size=9, color=LIGHT_BLUE)
    # Zurich logo top right
    add_zurich_logo(slide, inches(11.5), inches(0.15), w=inches(1.6), h=inches(0.55))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — EXECUTIVE SUMMARY (ONE-PAGER)
# ══════════════════════════════════════════════════════════════════════════════
slide_layout = prs.slide_layouts[6]  # blank
slide1 = prs.slides.add_slide(slide_layout)

# Full blue background header
hdr_h = inches(1.1)
add_rect(slide1, 0, 0, W, hdr_h, fill_color=ZURICH_BLUE)
add_textbox(slide1, inches(0.4), inches(0.15), inches(9), inches(0.55),
            "K9 Agent — DA Direkt Dog Insurance",
            font_size=26, bold=True, color=WHITE, font_name="Calibri")
add_textbox(slide1, inches(0.4), inches(0.68), inches(9), inches(0.3),
            "AI-native Research, Quote & Bind  ·  Zurich Agentic AI Hyper Challenge 2026",
            font_size=10, color=LIGHT_BLUE)
add_zurich_logo(slide1, inches(11.5), inches(0.25), w=inches(1.6), h=inches(0.6))

# Row 1: Team / Use Case / Learnings
row1_y = hdr_h + inches(0.15)
row1_h = inches(0.95)
cw = inches(4.2)
add_label_box(slide1, inches(0.3), row1_y, cw, row1_h,
              "Name of the Team",
              "Retail_Agentic_Commerce_LiviuGhenghea")
add_label_box(slide1, inches(4.6), row1_y, cw, row1_h,
              "Selected Use Case",
              "01 — Agentic Commerce for the Personal Agent Era")
add_label_box(slide1, inches(8.9), row1_y, inches(4.1), row1_h,
              "Main Learning from Experience",
              "Building directly on the Anthropic API (no framework) kept the system transparent and fast. Agent-to-agent commerce is feasible today.")

# Row 2: Team members
row2_y = row1_y + row1_h + inches(0.1)
add_label_box(slide1, inches(0.3), row2_y, W - inches(0.6), inches(0.6),
              "Team Members",
              "Liviu Ghenghea  ·  liviu.ghenghea@zurich.com")

# Approach
row3_y = row2_y + inches(0.7)
add_label_box(slide1, inches(0.3), row3_y, W - inches(0.6), inches(1.1),
              "How your approach is solving the problem",
              "K9 Agent replaces the existing linear web funnel with a fully AI-native, conversational insurance journey. "
              "A Zurich-hosted AI agent (Claude Sonnet) understands natural language, calls 12 live DA Direkt APIs in real time, "
              "and guides customers from first question to bound policy — without a single form or dropdown. "
              "Built for agent-to-agent commerce: a personal AI can interact directly with the Zurich K9 Agent, enabling autonomous quote and bind. "
              "Addresses all 5 AS-IS pain points: relevance, understanding, engagement, transparency, purchase guidance.")

# Results + Next Steps side by side
row4_y = row3_y + inches(1.2)
half_w = (W - inches(0.8)) / 2
add_label_box(slide1, inches(0.3), row4_y, half_w, inches(1.55),
              "What the results of the prototype are",
              "• Live prototype: https://k9agent.streamlit.app\n"
              "• Full journey: breed lookup → live pricing → lead creation → bind\n"
              "• 18/18 API smoke tests passing against DA Direkt beta\n"
              "• Real leads created with valid reference numbers\n"
              "• Handles natural language, follow-up questions, edge cases\n"
              "• Cost per full journey: ~€0.05 (Sonnet + Haiku)")
add_label_box(slide1, inches(0.3) + half_w + inches(0.2), row4_y, half_w, inches(1.55),
              "What are the next steps for scaling your solution",
              "• MCP server: expose 12 tools — any personal agent can call the Zurich Agent directly\n"
              "• Production APIs: swap beta URL → production, rotate API key\n"
              "• Multi-channel: same agent core serves WhatsApp, web chat, voice\n"
              "• Multi-product: cat insurance, liability, life with minimal changes\n"
              "• Prompt caching: reduce Sonnet costs ~80% at scale\n"
              "• Session persistence: Redis for conversation state across channels")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM APPROACH & SOLUTION DESIGN
# ══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(slide_layout)
slide_header(slide2, "1", "Problem Approach & Solution Design",
             "Replacing the linear web funnel with an AI-native, agent-to-agent insurance journey")

# Pain points label
add_section_label(slide2, inches(0.3), inches(1.0), inches(12), "THE PROBLEM — 5 AS-IS PAIN POINTS")

pain_y = inches(1.3)
pain_h = inches(1.25)
pw = inches(3.1)
pains = [
    ("Relevance", "Customers can't see how the product connects to their pet's specific needs. Answers don't match their actual questions."),
    ("Understanding", "Complex terminology, unclear plan differences, no breed/age guidance. Surgery vs full coverage not intuitive."),
    ("Engagement", "Customers can't visualise what claims look like in practice. Network coverage and real-life usage unclear."),
    ("Transparency + Drop-off", "Premium logic opaque. Pre-existing conditions cause drop-offs. Application feels complex and error-prone."),
]
for i, (title, body) in enumerate(pains):
    add_label_box(slide2, inches(0.3) + i * (pw + inches(0.1)), pain_y, pw, pain_h,
                  title, body, bg=LIGHT_BLUE, label_color=MID_BLUE)

# Solution label
add_section_label(slide2, inches(0.3), inches(2.65), inches(12), "THE SOLUTION — HOW K9 AGENT SOLVES IT")

sol_y = inches(2.95)
sol_h = inches(1.1)
sw = inches(4.1)
solutions = [
    ("AI-native conversation", "Natural language replaces forms. Claude understands 'my 2-year-old Golden Retriever named Luna' in one message. No forms, no dropdowns."),
    ("Live pricing + transparency", "Real-time prices across all 6 plans from DA Direkt APIs. Explains GOT factors, deductibles, exclusions and waiting periods in plain language."),
    ("Agent-to-agent commerce", "Personal AI agents can call the Zurich K9 Agent directly via tool-use protocol. Autonomous quote and bind — no human in the loop."),
]
for i, (title, body) in enumerate(solutions):
    add_label_box(slide2, inches(0.3) + i * (sw + inches(0.15)), sol_y, sw, sol_h,
                  title, body, bg=MID_BLUE, label_color=WHITE)

# Architecture line
arch_y = inches(4.15)
add_rect(slide2, inches(0.3), arch_y, W - inches(0.6), inches(0.22), fill_color=MID_BLUE)
add_section_label(slide2, inches(0.3), arch_y - inches(0.25), W - inches(0.6), "ARCHITECTURE IN ONE LINE")
add_textbox(slide2, inches(0.5), arch_y + inches(0.02), W - inches(1.0), inches(0.2),
            "Customer Agent (Haiku — fast, cheap)  →  natural language  →  Zurich K9 Agent (Sonnet — reasoning + 12 tools)  →  REST calls  →  DA Direkt Petolo APIs (live breeds / pricing / leads / bind)",
            font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Journey flow
flow_y = inches(4.55)
add_section_label(slide2, inches(0.3), flow_y, W - inches(0.6), "JOURNEY FLOW (non-linear)")
flow_y2 = flow_y + inches(0.25)
fw = inches(2.1)
fh = inches(0.95)
steps = [
    ("1  RESEARCH",  "Customer asks about\nplans & coverage"),
    ("2  QUOTE",     "Live prices fetched\nacross all 6 plans"),
    ("3  COLLECT",   "Owner + dog details\n+ IBAN gathered"),
    ("4  VALIDATE",  "check_missing_fields\nverifies readiness"),
    ("5  CONFIRM",   "Summary shown,\nexplicit consent"),
    ("6  BIND",      "Application submitted\nto DA Direkt ✓"),
]
colors_flow = [MID_BLUE, MID_BLUE, MID_BLUE, MID_BLUE, MID_BLUE, GREEN]
for i, (title, body) in enumerate(steps):
    fc = colors_flow[i]
    add_rect(slide2, inches(0.3) + i * (fw + inches(0.06)), flow_y2, fw, inches(0.3), fill_color=fc)
    add_textbox(slide2, inches(0.3) + i * (fw + inches(0.06)), flow_y2 + inches(0.04), fw, inches(0.25),
                title, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    body_bg = LIGHT_GREEN if fc == GREEN else LIGHT_BLUE
    add_rect(slide2, inches(0.3) + i * (fw + inches(0.06)), flow_y2 + inches(0.3), fw, fh - inches(0.3),
             fill_color=body_bg, line_color=fc, line_width=Pt(0.5))
    add_textbox(slide2, inches(0.35) + i * (fw + inches(0.06)), flow_y2 + inches(0.34), fw - inches(0.1), fh - inches(0.38),
                body, font_size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESULTS & KEY INSIGHTS
# ══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(slide_layout)
slide_header(slide3, "2", "Results & Key Insights",
             "Working prototype validated against live DA Direkt beta APIs")

# Metrics
add_section_label(slide3, inches(0.3), inches(1.0), W - inches(0.6), "KEY METRICS")
metrics = [
    ("18/18", "API Tests", "All smoke tests passing", GREEN),
    ("12", "Tools", "Live DA Direkt API calls", MID_BLUE),
    ("6", "Plans Quoted", "Real-time pricing all tiers", ZURICH_BLUE),
    ("~€0.05", "Cost/Journey", "Sonnet + Haiku combined", ORANGE),
    ("100%", "Conversational", "No forms, no dropdowns", GREEN),
]
mw = inches(2.4)
mh = inches(1.0)
for i, (val, lbl, sub, col) in enumerate(metrics):
    mx = inches(0.3) + i * (mw + inches(0.1))
    add_rect(slide3, mx, inches(1.28), mw, mh, fill_color=col)
    add_textbox(slide3, mx, inches(1.32), mw, inches(0.42),
                val, font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide3, mx, inches(1.73), mw, inches(0.22),
                lbl, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide3, mx, inches(1.93), mw, inches(0.22),
                sub, font_size=7, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Key Insights
add_section_label(slide3, inches(0.3), inches(2.45), W - inches(0.6), "KEY INSIGHTS")
insights = [
    ("Full journey end-to-end",
     "From first message to insurance application submission — live leads created in DA Direkt's beta system with real reference numbers. "
     "Deployed at https://k9agent.streamlit.app. Handles natural language, follow-ups, and edge cases without breaking."),
    ("Genuine AI understanding",
     "Agent understands natural language without rigid step sequences. Handles breed ambiguity, follow-up questions, and unexpected inputs. "
     "The conversation adapts to the customer — not the other way around."),
    ("Agent-to-agent ready",
     "Architecture designed so any MCP-compatible personal agent can discover and call the Zurich K9 Agent tools directly. "
     "No UI required — pure agent-to-agent API interaction with no human in the loop."),
]
iw = inches(4.2)
ih = inches(1.15)
for i, (title, body) in enumerate(insights):
    add_label_box(slide3, inches(0.3) + i * (iw + inches(0.15)), inches(2.72), iw, ih,
                  title, body, bg=MID_BLUE, label_color=WHITE)

# AS-IS vs TO-BE table
add_section_label(slide3, inches(0.3), inches(4.0), W - inches(0.6), "AS-IS vs TO-BE COMPARISON")
table_y = inches(4.28)
table_h = inches(2.8)
table_w = W - inches(0.6)
tbl = slide3.shapes.add_table(6, 3, inches(0.3), table_y, table_w, table_h).table

# Column widths
tbl.columns[0].width = inches(2.6)
tbl.columns[1].width = inches(4.5)
tbl.columns[2].width = inches(5.8)

headers = ["Dimension", "AS-IS (Today)", "K9 Agent (TO-BE)"]
rows_data = [
    ("Customer effort",         "Navigate 5-step web form",       "Single conversation, any device"),
    ("Coverage clarity",        "PDF documents, buried terms",     "Plain-language explanation on demand"),
    ("Pricing transparency",    "Calculator with fixed inputs",    "Live prices across all plans, scenario-based"),
    ("Pre-existing conditions", "Drop-off, unclear guidance",      "Proactive FAQ, guided disclosure"),
    ("Agent-to-agent",          "Not supported",                   "Native — personal agents can call directly"),
]

for col_i, hdr in enumerate(headers):
    cell = tbl.cell(0, col_i)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = ZURICH_BLUE
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

for row_i, (dim, asis, tobe) in enumerate(rows_data):
    bg = LIGHT_BLUE if row_i % 2 == 0 else WHITE
    for col_i, val in enumerate([dim, asis, tobe]):
        cell = tbl.cell(row_i + 1, col_i)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(8.5)
        run.font.name = "Calibri"
        if col_i == 2:
            run.font.color.rgb = GREEN
            run.font.bold = True
        else:
            run.font.color.rgb = DARK_GRAY


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NEXT STEPS & RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(slide_layout)
slide_header(slide4, "3", "Next Steps & Recommendations",
             "From hackathon prototype to enterprise-grade agentic distribution")

# Roadmap
add_section_label(slide4, inches(0.3), inches(1.0), W - inches(0.6), "ROADMAP")
roadmap = [
    ("NOW\n0–1 mo",  MID_BLUE,    "Production Readiness",
     "• Swap beta URL → production Petolo API\n"
     "• Rotate API key to production credential\n"
     "• Add Redis for session state persistence\n"
     "• Implement customer identity verification\n"
     "• Subscribe to policy issuance webhooks"),
    ("NEXT\n1–3 mo", ZURICH_BLUE, "Scale & Multi-channel",
     "• Expose tools as MCP server — any personal agent can discover and call\n"
     "• Add WhatsApp / voice channel wrappers\n"
     "• Extend to cat insurance and liability products\n"
     "• Enable prompt caching (80% cost reduction at scale)\n"
     "• Add conversation analytics dashboard"),
    ("FUTURE\n3–12 mo", GREEN,    "Enterprise Distribution",
     "• Open MCP server to approved personal agent platforms\n"
     "• B2B: white-label the Zurich Agent for broker networks\n"
     "• Cross-sell signals: agent detects life events, triggers relevant products\n"
     "• Claims co-pilot: extend same architecture to claims journey\n"
     "• Group rollout: replicate for other DA Direkt / Zurich markets"),
]
rw = inches(4.1)
rh = inches(2.4)
for i, (phase, col, title, bullets) in enumerate(roadmap):
    rx = inches(0.3) + i * (rw + inches(0.15))
    ry = inches(1.28)
    # Phase badge
    add_rect(slide4, rx, ry, inches(0.9), rh, fill_color=col)
    add_textbox(slide4, rx + inches(0.02), ry + inches(0.5), inches(0.86), inches(0.8),
                phase, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title bar
    add_rect(slide4, rx + inches(0.9), ry, rw - inches(0.9), inches(0.35), fill_color=col)
    add_textbox(slide4, rx + inches(0.95), ry + inches(0.05), rw - inches(1.0), inches(0.28),
                title, font_size=10, bold=True, color=WHITE)
    # Bullets
    add_rect(slide4, rx + inches(0.9), ry + inches(0.35), rw - inches(0.9), rh - inches(0.35),
             fill_color=WHITE, line_color=col, line_width=Pt(0.5))
    add_textbox(slide4, rx + inches(0.95), ry + inches(0.42), rw - inches(1.05), rh - inches(0.5),
                bullets, font_size=8, color=DARK_GRAY, wrap=True)

# Cost model
cost_y = inches(3.85)
add_section_label(slide4, inches(0.3), cost_y, W - inches(0.6), "COST MODEL")
cost_tbl = slide4.shapes.add_table(4, 5, inches(0.3), cost_y + inches(0.28), W - inches(0.6), inches(1.2)).table
cost_tbl.columns[0].width = inches(2.5)
cost_tbl.columns[1].width = inches(3.0)
cost_tbl.columns[2].width = inches(2.2)
cost_tbl.columns[3].width = inches(2.8)
cost_tbl.columns[4].width = inches(2.5)

cost_headers = ["Stage", "Model", "Cost per journey", "At 10k journeys/mo", "With prompt caching"]
cost_rows = [
    ("Zurich Agent",   "claude-sonnet-4-6", "~€0.05",   "~€500/mo",   "~€100/mo"),
    ("Customer Agent", "claude-haiku-4-5",  "~€0.003",  "~€30/mo",    "~€6/mo"),
    ("Total",          "—",                 "~€0.05",   "~€530/mo",   "~€106/mo"),
]

for ci, hdr in enumerate(cost_headers):
    cell = cost_tbl.cell(0, ci)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = ZURICH_BLUE
    p = cell.text_frame.paragraphs[0]
    run = p.runs[0]
    run.font.bold = True
    run.font.size = Pt(8.5)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

for ri, row in enumerate(cost_rows):
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
    if ri == 2:
        bg = LIGHT_GREEN
    for ci, val in enumerate(row):
        cell = cost_tbl.cell(ri + 1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(8.5)
        run.font.name = "Calibri"
        run.font.bold = (ri == 2)
        run.font.color.rgb = GREEN if ri == 2 else DARK_GRAY

# Vision box
vis_y = inches(5.15)
add_rect(slide4, inches(0.3), vis_y, W - inches(0.6), inches(0.82), fill_color=ZURICH_BLUE)
add_textbox(slide4, inches(0.5), vis_y + inches(0.08), W - inches(1.0), inches(0.65),
            "Vision:  The K9 Agent architecture is the blueprint for Zurich's agentic distribution layer. "
            "The same pattern — Zurich Agent + tool layer + product APIs — can power any insurance product, "
            "any channel, and any personal agent platform that adopts MCP.  Built once, distributed everywhere.",
            font_size=9, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"PPTX created: {OUTPUT}")
